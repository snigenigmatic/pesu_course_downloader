#!/usr/bin/env python3
"""
Interactive PESU Academy Downloader
Complete workflow: Login → Select Course → Select Units → Select Resources → Download → Convert → Merge
"""

import os
import sys
import json
import getpass
import shutil
from pathlib import Path
from typing import List, Dict, Optional, Tuple
import requests
from bs4 import BeautifulSoup
import re
import curses
from colorama import Fore, Style, init as colorama_init
from pypdf import PdfWriter
import subprocess
import tempfile
import time
import zipfile
import logging

# Try importing conversion libraries
try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    import comtypes.client
    COMTYPES_AVAILABLE = True
except ImportError:
    COMTYPES_AVAILABLE = False

try:
    import aspose.slides as slides
    ASPOSE_AVAILABLE = True
except ImportError:
    ASPOSE_AVAILABLE = False

# Initialize colorama
colorama_init(autoreset=True)

# Cache file for course data
CACHE_FILE = Path("courses.json")

# Resource type mapping (id parameter values)
RESOURCE_TYPES = {
    "2": "Slides",
    "3": "Notes",
    "4": "QA",
    "5": "Assignments",
    "6": "QB",
    "7": "MCQs",
    "8": "References",
}


class Timer:
    """Simple timer for performance tracking"""
    def __init__(self):
        self.start = time.time()
    
    def elapsed(self) -> float:
        return time.time() - self.start
    
    def pretty(self) -> str:
        elapsed = self.elapsed()
        if elapsed < 1:
            return f"{elapsed*1000:.0f}ms"
        return f"{elapsed:.1f}s"


class Spinner:
    """Context manager for showing a spinner during operations"""
    def __init__(self, message: str):
        self.message = message
    
    def __enter__(self):
        print(f"  {self.message}...", end="", flush=True)
        return self
    
    def __exit__(self, *args):
        print("")


class DOCXRepair:
    """Handles various DOCX repair strategies"""
    
    def __init__(self):
        pass
    
    def repair_with_docx(self, input_path: Path, output_path: Path) -> bool:
        """Repair by loading and re-saving with python-docx"""
        try:
            import docx
            doc = docx.Document(str(input_path))
            doc.save(str(output_path))
            return True
        except Exception:
            return False
    
    def repair_by_rezip(self, input_path: Path, output_path: Path) -> bool:
        """Repair by extracting and re-zipping"""
        try:
            with tempfile.TemporaryDirectory() as temp_dir:
                temp_path = Path(temp_dir)
                try:
                    with zipfile.ZipFile(input_path, 'r') as zip_ref:
                        zip_ref.extractall(temp_path)
                except zipfile.BadZipFile:
                    return False
                
                with zipfile.ZipFile(output_path, 'w', zipfile.ZIP_DEFLATED) as zip_out:
                    for root, dirs, files in os.walk(temp_path):
                        for file in files:
                            file_path = Path(root) / file
                            arcname = file_path.relative_to(temp_path)
                            zip_out.write(file_path, arcname)
                return True
        except Exception:
            return False
    
    def repair_xml_relationships(self, input_path: Path, output_path: Path) -> bool:
        """Repair broken document XML relationships"""
        try:
            with tempfile.TemporaryDirectory() as temp_dir:
                temp_path = Path(temp_dir)
                try:
                    with zipfile.ZipFile(input_path, 'r') as zip_ref:
                        zip_ref.extractall(temp_path)
                except zipfile.BadZipFile:
                    return False
                
                # Fix relationships in word/_rels
                rels_dir = temp_path / "word" / "_rels"
                if rels_dir.exists():
                    for rels_file in rels_dir.glob("*.rels"):
                        try:
                            content = rels_file.read_text(encoding='utf-8')
                            # Remove empty hyperlinks
                            content = content.replace('Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink" TargetMode="External" Target=""', '')
                            # Remove other problematic patterns
                            content = re.sub(r'<Relationship[^>]*Target=""[^>]*/>', '', content)
                            rels_file.write_text(content, encoding='utf-8')
                        except Exception:
                            pass
                
                # Fix document.xml if needed
                doc_xml = temp_path / "word" / "document.xml"
                if doc_xml.exists():
                    try:
                        content = doc_xml.read_text(encoding='utf-8')
                        # Remove invalid XML characters
                        content = re.sub(r'[\x00-\x08\x0B\x0C\x0E-\x1F]', '', content)
                        doc_xml.write_text(content, encoding='utf-8')
                    except Exception:
                        pass
                
                with zipfile.ZipFile(output_path, 'w', zipfile.ZIP_DEFLATED) as zip_out:
                    for root, dirs, files in os.walk(temp_path):
                        for file in files:
                            file_path = Path(root) / file
                            arcname = file_path.relative_to(temp_path)
                            zip_out.write(file_path, arcname)
                
                # Verify the repair worked
                try:
                    import docx
                    docx.Document(str(output_path))
                    return True
                except:
                    # If python-docx not available, just check if file is valid
                    return True
        except Exception:
            return False
    
    def attempt_repair(self, input_path: Path) -> Optional[Path]:
        """Attempt all repair strategies in order"""
        temp_dir = Path(tempfile.mkdtemp())
        strategies = [
            self.repair_by_rezip,
            self.repair_xml_relationships,
            self.repair_with_docx,
        ]
        
        for i, strategy in enumerate(strategies):
            output_path = temp_dir / f"repaired_{i}_{input_path.name}"
            if strategy(input_path, output_path):
                if output_path.exists() and output_path.stat().st_size > 0:
                    return output_path
        
        shutil.rmtree(temp_dir, ignore_errors=True)
        return None


class PPTXRepair:
    """Handles various PPTX repair strategies"""
    
    def __init__(self):
        pass
    
    def repair_with_pptx(self, input_path: Path, output_path: Path) -> bool:
        """Repair by loading and re-saving with python-pptx"""
        if not PPTX_AVAILABLE:
            return False
        try:
            prs = Presentation(str(input_path))
            prs.save(str(output_path))
            return True
        except Exception:
            return False
    
    def repair_by_rezip(self, input_path: Path, output_path: Path) -> bool:
        """Repair by extracting and re-zipping"""
        try:
            with tempfile.TemporaryDirectory() as temp_dir:
                temp_path = Path(temp_dir)
                try:
                    with zipfile.ZipFile(input_path, 'r') as zip_ref:
                        zip_ref.extractall(temp_path)
                except zipfile.BadZipFile:
                    return False
                
                with zipfile.ZipFile(output_path, 'w', zipfile.ZIP_DEFLATED) as zip_out:
                    for root, dirs, files in os.walk(temp_path):
                        for file in files:
                            file_path = Path(root) / file
                            arcname = file_path.relative_to(temp_path)
                            zip_out.write(file_path, arcname)
                return True
        except Exception:
            return False
    
    def repair_xml_relationships(self, input_path: Path, output_path: Path) -> bool:
        """Repair broken slide XML relationships"""
        try:
            with tempfile.TemporaryDirectory() as temp_dir:
                temp_path = Path(temp_dir)
                try:
                    with zipfile.ZipFile(input_path, 'r') as zip_ref:
                        zip_ref.extractall(temp_path)
                except zipfile.BadZipFile:
                    return False
                
                rels_dir = temp_path / "ppt" / "_rels"
                if rels_dir.exists():
                    for rels_file in rels_dir.glob("*.rels"):
                        try:
                            content = rels_file.read_text(encoding='utf-8')
                            content = content.replace('Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink" TargetMode="External" Target=""', '')
                            rels_file.write_text(content, encoding='utf-8')
                        except Exception:
                            pass
                
                with zipfile.ZipFile(output_path, 'w', zipfile.ZIP_DEFLATED) as zip_out:
                    for root, dirs, files in os.walk(temp_path):
                        for file in files:
                            file_path = Path(root) / file
                            arcname = file_path.relative_to(temp_path)
                            zip_out.write(file_path, arcname)
                
                if PPTX_AVAILABLE:
                    try:
                        Presentation(str(output_path))
                        return True
                    except:
                        return False
                return True
        except Exception:
            return False
    
    def repair_by_rezip_linux(self, input_path: Path, output_path: Path) -> bool:
        """Linux-specific repair: extract and repackage using zipfile (mimics unzip+zip CLI)"""
        try:
            with tempfile.TemporaryDirectory() as temp_dir:
                temp_path = Path(temp_dir)
                try:
                    with zipfile.ZipFile(input_path, 'r') as zip_ref:
                        zip_ref.extractall(temp_path)
                except zipfile.BadZipFile:
                    return False

                with zipfile.ZipFile(output_path, 'w', zipfile.ZIP_DEFLATED) as zip_out:
                    for root, dirs, files in os.walk(temp_path):
                        for file in sorted(files):  # sorted for deterministic order
                            file_path = Path(root) / file
                            arcname = file_path.relative_to(temp_path)
                            zip_out.write(file_path, arcname)
                return output_path.exists() and output_path.stat().st_size > 0
        except Exception:
            return False
    
    def attempt_repair(self, input_path: Path) -> Optional[Path]:
        """Attempt all repair strategies in order"""
        temp_dir = Path(tempfile.mkdtemp())
        strategies = [
            self.repair_with_pptx,
            self.repair_by_rezip,
            self.repair_xml_relationships
        ]
        
        for i, strategy in enumerate(strategies):
            output_path = temp_dir / f"repaired_{i}_{input_path.name}"
            if strategy(input_path, output_path):
                if output_path.exists() and output_path.stat().st_size > 0:
                    return output_path
        
        shutil.rmtree(temp_dir, ignore_errors=True)
        return None


class OfficeConverter:
    """Handles Office file to PDF conversion with multiple methods"""
    
    def __init__(self):
        self.pptx_repairer = PPTXRepair()
        self.docx_repairer = DOCXRepair()
    
    def convert_with_powerpoint(self, input_path: Path, output_path: Path) -> bool:
        """Convert using Microsoft PowerPoint via COM automation"""
        if not COMTYPES_AVAILABLE:
            return False
        
        try:
            ppSaveAsPDF = 32
            powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
            powerpoint.Visible = 1
            
            try:
                deck = powerpoint.Presentations.Open(str(input_path.absolute()), WithWindow=False)
                deck.SaveAs(str(output_path.absolute()), ppSaveAsPDF)
                deck.Close()
                return True
            finally:
                powerpoint.Quit()
        except Exception:
            return False
    
    def convert_with_word(self, input_path: Path, output_path: Path) -> bool:
        """Convert using Microsoft Word via COM automation"""
        if not COMTYPES_AVAILABLE:
            return False
        
        try:
            wdFormatPDF = 17
            word = comtypes.client.CreateObject("Word.Application")
            word.Visible = False
            
            try:
                doc = word.Documents.Open(str(input_path.absolute()))
                doc.SaveAs(str(output_path.absolute()), FileFormat=wdFormatPDF)
                doc.Close()
                return True
            finally:
                word.Quit()
        except Exception:
            return False
    
    def convert_with_aspose_slides(self, input_path: Path, output_path: Path) -> bool:
        """Convert PPTX using Aspose.Slides"""
        if not ASPOSE_AVAILABLE:
            return False
        try:
            presentation = slides.Presentation(str(input_path))
            presentation.save(str(output_path), slides.export.SaveFormat.PDF)
            return True
        except Exception:
            return False
    
    def convert_with_libreoffice(self, input_path: Path, output_path: Path) -> bool:
        """Convert using LibreOffice (cross-platform)"""
        import platform
        is_linux = platform.system() == "Linux"

        if is_linux:
            # On Linux, find soffice in PATH
            soffice = shutil.which("soffice") or shutil.which("libreoffice")
            if not soffice:
                return False

            try:
                subprocess.run(
                    [soffice, "--headless", "--convert-to", "pdf",
                     "--outdir", str(input_path.parent), str(input_path)],
                    capture_output=True, text=True, timeout=120
                )
                # LibreOffice creates PDF next to input file
                generated = input_path.with_suffix(".pdf")
                if generated.exists():
                    shutil.move(str(generated), str(output_path))  # safely move across filesystems
                    return output_path.exists() and output_path.stat().st_size > 0
                return False
            except Exception:
                return False

        else:
            # Original Windows logic unchanged
            possible_paths = [
                r"C:\Program Files\LibreOffice\program\soffice.exe",
                r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
                r"C:\Program Files\LibreOffice 7\program\soffice.exe",
                r"C:\Program Files\LibreOffice 24\program\soffice.exe",
            ]
            soffice = None
            for path in possible_paths:
                if Path(path).exists():
                    soffice = path
                    break
            if not soffice:
                return False

            for process_name in ["soffice.exe", "soffice.bin"]:
                try:
                    subprocess.run(["taskkill", "/F", "/IM", process_name, "/T"],
                                   capture_output=True, timeout=3)
                except:
                    pass
            time.sleep(0.5)

            try:
                if output_path.exists():
                    output_path.unlink()
                cmd = [soffice, "--headless", "--convert-to", "pdf",
                       "--outdir", str(output_path.parent), str(input_path)]
                result = subprocess.run(
                    cmd, capture_output=True, text=True, timeout=120,
                    creationflags=subprocess.CREATE_NO_WINDOW if hasattr(subprocess, 'CREATE_NO_WINDOW') else 0
                )
                time.sleep(2)
                return output_path.exists() and output_path.stat().st_size > 0
            except Exception:
                return False
            finally:
                for process_name in ["soffice.exe", "soffice.bin"]:
                    try:
                        subprocess.run(["taskkill", "/F", "/IM", process_name, "/T"],
                                       capture_output=True, timeout=3)
                    except:
                        pass
    
    def convert_pptx_to_pdf(self, input_path: Path, output_path: Path) -> Tuple[bool, str]:
        """Convert PPTX to PDF with repair and multiple conversion methods"""
        import platform
        is_linux = platform.system() == "Linux"

        if is_linux:
            # On Linux: try repair-then-libreoffice first (most reliable)
            repaired_path = None

            # Try python-pptx repair first
            temp_dir = Path(tempfile.mkdtemp())
            repaired_candidate = temp_dir / f"repaired_{input_path.name}"

            if self.pptx_repairer.repair_by_rezip_linux(input_path, repaired_candidate):
                repaired_path = repaired_candidate

            target = repaired_path if repaired_path else input_path
            method_suffix = " (repaired)" if repaired_path else ""

            if self.convert_with_libreoffice(target, output_path):
                if output_path.exists() and output_path.stat().st_size > 0:
                    try:
                        shutil.rmtree(temp_dir, ignore_errors=True)
                    except:
                        pass
                    return True, f"LibreOffice{method_suffix}"

            # Fallback: Aspose if available
            if self.convert_with_aspose_slides(target, output_path):
                if output_path.exists() and output_path.stat().st_size > 0:
                    return True, f"Aspose.Slides{method_suffix}"

            try:
                shutil.rmtree(temp_dir, ignore_errors=True)
            except:
                pass
            return False, "none"

        else:
            # Original Windows logic unchanged
            if self.convert_with_powerpoint(input_path, output_path):
                if output_path.exists() and output_path.stat().st_size > 0:
                    return True, "PowerPoint COM"

            repaired_path = self.pptx_repairer.attempt_repair(input_path)
            if repaired_path:
                methods = [
                    (self.convert_with_powerpoint, "PowerPoint COM"),
                    (self.convert_with_libreoffice, "LibreOffice"),
                    (self.convert_with_aspose_slides, "Aspose.Slides")
                ]
                for method, method_name in methods:
                    if method(repaired_path, output_path):
                        if output_path.exists() and output_path.stat().st_size > 0:
                            try:
                                shutil.rmtree(repaired_path.parent, ignore_errors=True)
                            except:
                                pass
                            return True, f"{method_name} (repaired)"
                try:
                    shutil.rmtree(repaired_path.parent, ignore_errors=True)
                except:
                    pass

            for method, method_name in [
                (self.convert_with_aspose_slides, "Aspose.Slides"),
                (self.convert_with_libreoffice, "LibreOffice")
            ]:
                if method(input_path, output_path):
                    if output_path.exists() and output_path.stat().st_size > 0:
                        return True, method_name

            return False, "none"
    
    def convert_docx_to_pdf(self, input_path: Path, output_path: Path) -> Tuple[bool, str]:
        """Convert DOCX to PDF with repair and multiple conversion methods"""
        # Try Word COM first (best quality)
        if self.convert_with_word(input_path, output_path):
            if output_path.exists() and output_path.stat().st_size > 0:
                return True, "Word COM"
        
        # Try repair strategies if direct conversion failed
        repaired_path = self.docx_repairer.attempt_repair(input_path)
        if repaired_path:
            methods = [
                (self.convert_with_word, "Word COM"),
                (self.convert_with_libreoffice, "LibreOffice")
            ]
            for method, method_name in methods:
                if method(repaired_path, output_path):
                    if output_path.exists() and output_path.stat().st_size > 0:
                        # Cleanup repaired file
                        try:
                            shutil.rmtree(repaired_path.parent, ignore_errors=True)
                        except:
                            pass
                        return True, f"{method_name} (repaired)"
            
            # Cleanup failed repair
            try:
                shutil.rmtree(repaired_path.parent, ignore_errors=True)
            except:
                pass
        
        # If repair failed, try LibreOffice as last resort
        if self.convert_with_libreoffice(input_path, output_path):
            if output_path.exists() and output_path.stat().st_size > 0:
                return True, "LibreOffice"
        
        return False, "none"


class PESUInteractiveDownloader:
    def __init__(self, username: str, password: str):
        self.session = requests.Session()
        self.username = username
        self.password = password
        self.base_url = "https://www.pesuacademy.com/Academy"
        self.downloaded_files = []

    def detect_file_type(self, file_path: Path) -> Optional[str]:
        """Detect actual file type from magic bytes and content structure"""
        try:
            with open(file_path, "rb") as f:
                header = f.read(8)
            
            # Check magic bytes
            if header.startswith(b'PK\x03\x04'):
                # It's a ZIP file - could be Office format (DOCX/PPTX/XLSX)
                # Read more content to check for Office format markers
                with open(file_path, "rb") as f:
                    content = f.read(2048)  # Read more bytes for better detection
                    content_str = content.decode('latin-1', errors='ignore')
                    
                    # Check for specific Office document markers (order matters!)
                    # DOCX has word/ directory structure
                    if 'word/' in content_str or 'word/_rels' in content_str or '[Content_Types].xml' in content_str and 'wordprocessingml' in content_str:
                        return '.docx'
                    # PPTX has ppt/ directory structure
                    elif 'ppt/' in content_str or 'ppt/_rels' in content_str or 'slideshow' in content_str.lower() or 'presentationml' in content_str:
                        return '.pptx'
                    # XLSX has xl/ directory structure
                    elif 'xl/' in content_str or 'xl/_rels' in content_str or 'workbook' in content_str.lower() or 'spreadsheetml' in content_str:
                        return '.xlsx'
                    else:
                        # If we can't determine, try to inspect ZIP contents
                        try:
                            import zipfile
                            with zipfile.ZipFile(file_path, 'r') as zf:
                                namelist = zf.namelist()
                                if any('word/' in name for name in namelist):
                                    return '.docx'
                                elif any('ppt/' in name for name in namelist):
                                    return '.pptx'
                                elif any('xl/' in name for name in namelist):
                                    return '.xlsx'
                        except:
                            pass
                        # If still unknown, return None instead of guessing
                        return None
            elif header.startswith(b'%PDF'):
                return '.pdf'
            elif header.startswith(b'\xd0\xcf\x11\xe0'):
                # Old Office format (DOC/PPT/XLS) - compound file binary format
                # Try to distinguish by reading more
                with open(file_path, "rb") as f:
                    content = f.read(2048)
                    content_str = content.decode('latin-1', errors='ignore')
                    if 'Word' in content_str or 'Microsoft Word' in content_str:
                        return '.doc'
                    elif 'PowerPoint' in content_str or 'Microsoft PowerPoint' in content_str:
                        return '.ppt'
                    else:
                        return '.ppt'  # Default for old format
            
            return None
        except Exception:
            return None

    def login(self):
        """Login to PESU Academy"""
        print(f"\n{Fore.CYAN}[1/7] Logging in...{Style.RESET_ALL}")
        r0 = self.session.get(f"{self.base_url}/")
        soup = BeautifulSoup(r0.text, "html.parser")
        csrf_input = soup.find("input", {"name": "_csrf"})
        csrf_token = csrf_input.get("value") if csrf_input else None

        if not csrf_token:
            raise Exception("Could not find CSRF token")

        login_data = {
            "j_username": self.username,
            "j_password": self.password,
            "_csrf": csrf_token,
        }

        response = self.session.post(
            f"{self.base_url}/j_spring_security_check", data=login_data
        )

        if "authfailed" in response.url:
            raise Exception("Login failed! Check credentials.")

        print(f"{Fore.GREEN}✓ Login successful{Style.RESET_ALL}")

    def get_courses(self) -> List[Dict]:
        """Get all available courses (cache-first approach)"""
        print(f"\n{Fore.CYAN}[2/7] Fetching available courses...{Style.RESET_ALL}")

        # ── always try cache first ────────────────────────────────────
        if CACHE_FILE.exists():
            try:
                with open(CACHE_FILE, "r", encoding="utf-8") as f:
                    cached = json.load(f)
                courses = cached.get("courses", [])
                saved_at = cached.get("_saved_at", 0)
                age_str = self._cache_age_str(saved_at)
                if courses:
                    print(f"{Fore.GREEN}✓ Loaded {len(courses)} courses from cache  "
                          f"{Fore.YELLOW}(saved {age_str}){Style.RESET_ALL}")
                    return courses
            except Exception:
                print(f"{Fore.YELLOW}  Cache unreadable, fetching from server...{Style.RESET_ALL}")

        # ── no cache at all → must fetch ─────────────────────────────
        print(f"{Fore.YELLOW}  No cache found — fetching from server (first run, one-time wait)...{Style.RESET_ALL}")
        return self._fetch_and_cache_courses()

    def _cache_age_str(self, saved_at: float) -> str:
        """Get human-readable cache age"""
        age = time.time() - saved_at
        if age < 3600:   return f"{age/60:.0f}m ago"
        if age < 86400:  return f"{age/3600:.0f}h ago"
        return f"{age/86400:.0f}d ago"

    def _fetch_and_cache_courses(self) -> List[Dict]:
        """Fetch courses from server and cache them"""
        t = Timer()
        with Spinner("Fetching course list from server"):
            response = self.session.get(f"{self.base_url}/a/g/getSubjectsCode")

        if response.status_code != 200:
            raise Exception("Failed to fetch courses")

        soup = BeautifulSoup(response.text, "html.parser")
        courses = []
        for option in soup.find_all("option"):
            course_id = option.get("value")
            course_name = option.text.strip()
            if course_id and course_name:
                course_id = str(course_id).strip().replace('\\"', "").replace("\\'", "").strip('"').strip("'").replace("\\", "")
                subject_code = course_name.split("-")[0].strip() if "-" in course_name else course_name
                courses.append({"id": course_id, "subjectCode": subject_code, "subjectName": course_name})

        print(f"{Fore.GREEN}  Found {len(courses)} courses  {Fore.YELLOW}({t.pretty()}){Style.RESET_ALL}")

        try:
            with open(CACHE_FILE, "w", encoding="utf-8") as f:
                json.dump({"_saved_at": time.time(), "courses": courses}, f, indent=2, ensure_ascii=False)
            print(f"  {Fore.CYAN}Saved to courses.json{Style.RESET_ALL}")
        except Exception as e:
            print(f"  {Fore.YELLOW}⚠ Could not save cache: {e}{Style.RESET_ALL}")

        return courses

    def filter_courses_by_year(self, courses: List[Dict]) -> List[Dict]:
        """Filter and sort courses by academic year"""
        print(f"\n{Fore.CYAN}Select Academic Year:{Style.RESET_ALL}")
        print("  1. UE25")
        print("  2. UE24")
        print("  3. UE23")
        print("  4. UE22")
        print("  5. UE21")
        print("  6. UE20")
        print("  7. All years")
        
        choice = input(f"\n{Fore.CYAN}Enter choice (1-7, default=3): {Style.RESET_ALL}").strip() or "3"
        
        filter_map = {
            "1": ["UE25"],
            "2": ["UE24"],
            "3": ["UE23"],
            "4": ["UE22"],
            "5": ["UE21"],
            "6": ["UE20"],
            "7": ["UE25", "UE24", "UE23", "UE22", "UE21", "UE20"],
        }
        
        prefixes = filter_map.get(choice, ["UE23"])
        
        # Filter courses
        filtered = [
            c for c in courses
            if any(c["subjectCode"].startswith(prefix) for prefix in prefixes)
        ]
        
        # Sort by year (newest first) then alphabetically
        def get_year_priority(course):
            code = course["subjectCode"]
            for i, prefix in enumerate(["UE25", "UE24", "UE23", "UE22", "UE21", "UE20"]):
                if code.startswith(prefix):
                    return i
            return 99

        filtered.sort(key=lambda c: (get_year_priority(c), c["subjectCode"]))

        year_names = {
            "1": "2025-26",
            "2": "2024-25",
            "3": "2023-24",
            "4": "2022-23",
            "5": "2021-22",
            "6": "2020-21",
            "7": "recent years"
        }
        
        print(f"{Fore.GREEN}✓ Filtered to {len(filtered)} courses for {year_names.get(choice, 'selected years')}{Style.RESET_ALL}")
        return filtered

    def get_units(self, course_id: str) -> List[Dict]:
        """Get all units for a course"""
        url = f"{self.base_url}/a/i/getCourse/{course_id}"
        response = self.session.get(url)
        soup = BeautifulSoup(response.text, "html.parser")
        options = soup.find_all("option")

        units = []
        for option in options:
            unit_id = option.get("value")
            unit_name = option.text.strip()
            if unit_id and unit_name:
                unit_id = (
                    str(unit_id).strip().replace("\\", "").strip('"').strip("'")
                )
                units.append({"id": unit_id, "name": unit_name})

        return units

    def get_classes(self, unit_id: str) -> List[Dict]:
        """Get all classes for a unit"""
        url = f"{self.base_url}/a/i/getCourseClasses/{unit_id}"
        response = self.session.get(url)

        html_content = (
            response.json()
            if response.headers.get("Content-Type", "").startswith("application/json")
            else response.text
        )
        soup = BeautifulSoup(html_content, "html.parser")
        options = soup.find_all("option")

        classes = []
        for option in options:
            class_id = option.get("value")
            class_name = option.text.strip()
            if class_id and class_name:
                class_id = (
                    str(class_id).strip().replace("\\", "").strip('"').strip("'")
                )
                classes.append({"id": class_id, "name": class_name})

        return classes

    def get_resource_links(
        self, course_id: str, class_id: str, resource_type_id: str
    ) -> List[Dict]:
        """Get download links for a specific resource type"""
        url = f"{self.base_url}/s/studentProfilePESUAdmin"
        params = {
            "url": "studentProfilePESUAdmin",
            "controllerMode": "6403",
            "actionType": "60",
            "selectedData": course_id,
            "id": resource_type_id,
            "unitid": class_id,
        }

        response = self.session.get(url, params=params)

        # Check if direct file download
        content_type = response.headers.get("Content-Type", "")
        if "application/" in content_type and "html" not in content_type:
            return [
                {"type": "direct", "url": response.url, "content": response.content}
            ]

        # Parse HTML for download links
        soup = BeautifulSoup(response.text, "html.parser")
        download_links = []

        # Look for various download patterns
        for element in soup.find_all(onclick=True):
            onclick = element.get("onclick", "")
            text = element.text.strip()

            # Pattern 1: downloadslidecoursedoc in loadIframe
            if "downloadslidecoursedoc" in onclick:
                match = re.search(r"loadIframe\('([^']+)'", onclick)
                if match:
                    download_url = match.group(1).split("#")[0]
                    if download_url.startswith("/Academy"):
                        full_url = f"https://www.pesuacademy.com{download_url}"
                        download_links.append(
                            {"type": "link", "url": full_url, "text": text}
                        )

            # Pattern 2: downloadcoursedoc
            elif "downloadcoursedoc" in onclick:
                match = re.search(r"downloadcoursedoc\('([^']+)'", onclick)
                if match:
                    doc_id = match.group(1)
                    full_url = (
                        f"{self.base_url}/s/referenceMeterials/downloadcoursedoc/{doc_id}"
                    )
                    download_links.append(
                        {"type": "link", "url": full_url, "text": text}
                    )

        return download_links

    def get_filename_from_response(self, response) -> Optional[str]:
        """Extract filename from response headers"""
        content_disp = response.headers.get("Content-Disposition", "")
        if "filename=" in content_disp:
            match = re.search(
                r'filename[*]?=["\']?(?:UTF-8\'\')?([^"\';\n]+)', content_disp
            )
            if match:
                return match.group(1).strip()
        return None

    def download_resources(
        self,
        course_id: str,
        course_name: str,
        selected_units: List[int],
        selected_resources: List[str],
        base_dir: Path,
    ):
        """Download selected resources for selected units"""
        print(f"\n{Fore.CYAN}[4/7] Downloading resources...{Style.RESET_ALL}")

        units = self.get_units(course_id)
        total_downloaded = 0

        for unit_idx in selected_units:
            if unit_idx > len(units):
                print(f"{Fore.YELLOW}⚠ Unit {unit_idx} not found, skipping{Style.RESET_ALL}")
                continue

            unit = units[unit_idx - 1]
            print(f"\n{Fore.BLUE}{'='*70}{Style.RESET_ALL}")
            print(f"{Fore.BLUE}Unit {unit_idx}: {unit['name']}{Style.RESET_ALL}")
            print(f"{Fore.BLUE}{'='*70}{Style.RESET_ALL}")

            classes = self.get_classes(unit["id"])
            print(f"Found {len(classes)} classes")

            # Create unit directory
            unit_dir = base_dir / f"Unit_{unit_idx}"
            unit_dir.mkdir(parents=True, exist_ok=True)

            # Initialize counters for each resource type (separate counter per resource type)
            resource_counters = {res_id: 1 for res_id in selected_resources}

            # Process each class
            for class_idx, cls in enumerate(classes, 1):
                print(f"\n[{class_idx}/{len(classes)}] {cls['name']}")

                # Clean class name for filename
                safe_class_name = "".join(
                    c if c.isalnum() or c in (" ", "-", "_") else "_"
                    for c in cls["name"]
                ).strip()[:60]
                safe_class_name = "_".join(safe_class_name.split())

                # Try selected resource types
                for resource_id in selected_resources:
                    resource_name = RESOURCE_TYPES[resource_id]
                    links = self.get_resource_links(course_id, cls["id"], resource_id)

                    if links:
                        print(f"  {resource_name}: {len(links)} file(s)")

                        # Create resource type folder directly under unit (simplified structure)
                        resource_dir = unit_dir / resource_name
                        resource_dir.mkdir(parents=True, exist_ok=True)

                        # Download each file (use resource-specific counter)
                        for link in links:
                            current_counter = resource_counters[resource_id]
                            
                            if link["type"] == "direct":
                                # Direct download
                                filename = f"{current_counter}.{safe_class_name}.pdf"
                                output_path = resource_dir / filename
                                with open(output_path, "wb") as f:
                                    f.write(link["content"])
                                if output_path.stat().st_size > 0:
                                    # Detect actual file type from magic bytes
                                    actual_ext = self.detect_file_type(output_path)
                                    if actual_ext and actual_ext != '.pdf':
                                        # Rename with correct extension
                                        new_filename = f"{current_counter}.{safe_class_name}{actual_ext}"
                                        new_path = resource_dir / new_filename
                                        output_path.rename(new_path)
                                        output_path = new_path
                                        filename = new_filename
                                    
                                    print(
                                        f"    [OK] {filename} ({output_path.stat().st_size:,} bytes)"
                                    )
                                    total_downloaded += 1
                                    self.downloaded_files.append(output_path)
                                    resource_counters[resource_id] += 1  # Increment resource-specific counter
                            else:
                                # Download from link
                                headers = {
                                    "Referer": f"{self.base_url}/s/studentProfilePESU"
                                }
                                response = self.session.get(
                                    link["url"], headers=headers, stream=True
                                )

                                # Generate clean filename with sequential number
                                ext = ".pdf"
                                content_type = response.headers.get("Content-Type", "")
                                if "application/vnd.openxmlformats" in content_type:
                                    if "presentation" in content_type:
                                        ext = ".pptx"
                                    elif "word" in content_type:
                                        ext = ".docx"
                                    elif "sheet" in content_type:
                                        ext = ".xlsx"
                                elif "application/vnd.ms-powerpoint" in content_type:
                                    ext = ".ppt"
                                elif "application/msword" in content_type:
                                    ext = ".doc"
                                
                                # Clean filename: number.ClassName.ext (using resource-specific counter)
                                filename = f"{current_counter}.{safe_class_name}{ext}"

                                output_path = resource_dir / filename

                                # Save the file
                                with open(output_path, "wb") as f:
                                    for chunk in response.iter_content(chunk_size=8192):
                                        f.write(chunk)

                                if output_path.stat().st_size > 0:
                                    # Detect actual file type from magic bytes
                                    actual_ext = self.detect_file_type(output_path)
                                    if actual_ext and actual_ext != ext:
                                        # Rename with correct extension
                                        new_filename = f"{current_counter}.{safe_class_name}{actual_ext}"
                                        new_path = resource_dir / new_filename
                                        output_path.rename(new_path)
                                        output_path = new_path
                                        filename = new_filename
                                    
                                    print(
                                        f"    [OK] {filename} ({output_path.stat().st_size:,} bytes)"
                                    )
                                    total_downloaded += 1
                                    self.downloaded_files.append(output_path)
                                    resource_counters[resource_id] += 1  # Increment resource-specific counter
                                else:
                                    print(f"    [SKIP] {filename} (empty file, deleted)")
                                    output_path.unlink()

        print(f"\n{Fore.GREEN}{'='*70}{Style.RESET_ALL}")
        print(f"{Fore.GREEN}✓ Downloaded {total_downloaded} files{Style.RESET_ALL}")
        print(f"{Fore.GREEN}{'='*70}{Style.RESET_ALL}")


def convert_office_to_pdf(input_folder: Path) -> List[Path]:
    """Convert DOCX/PPTX files to PDF using advanced conversion methods"""
    print(f"\n{Fore.CYAN}[5/7] Converting files to PDF...{Style.RESET_ALL}")
    
    # Find all Office files first to determine what we're converting
    office_files = []
    for ext in ["*.docx", "*.pptx", "*.doc", "*.ppt"]:
        office_files.extend(input_folder.rglob(ext))
    
    if not office_files:
        print(f"{Fore.YELLOW}No Office files to convert{Style.RESET_ALL}")
        return []
    
    # Detect file types present
    has_pptx = any(f.suffix.lower() in ['.pptx', '.ppt'] for f in office_files)
    has_docx = any(f.suffix.lower() in ['.docx', '.doc'] for f in office_files)
    
    # Check for python-docx availability
    try:
        import docx
        DOCX_AVAILABLE = True
    except ImportError:
        DOCX_AVAILABLE = False
    
    # Show available conversion methods based on file types present
    print(f"\nAvailable conversion methods:")
    
    if has_pptx and has_docx:
        print(f"  PPTX files:")
        print(f"    • PowerPoint COM:  {'✓' if COMTYPES_AVAILABLE else '✗'} (Windows, best quality)")
        print(f"    • python-pptx:     {'✓' if PPTX_AVAILABLE else '✗'} (for repair)")
        print(f"    • Aspose.Slides:   {'✓' if ASPOSE_AVAILABLE else '✗'} (cross-platform)")
        print(f"  DOCX files:")
        print(f"    • Word COM:        {'✓' if COMTYPES_AVAILABLE else '✗'} (Windows, best quality)")
        print(f"    • python-docx:     {'✓' if DOCX_AVAILABLE else '✗'} (for repair)")
        print(f"  Both:")
        print(f"    • LibreOffice:     checking...")
    elif has_pptx:
        print(f"  • PowerPoint COM:  {'✓' if COMTYPES_AVAILABLE else '✗'} (Windows, best quality)")
        print(f"  • python-pptx:     {'✓' if PPTX_AVAILABLE else '✗'} (for repair)")
        print(f"  • Aspose.Slides:   {'✓' if ASPOSE_AVAILABLE else '✗'} (cross-platform)")
        print(f"  • LibreOffice:     checking...")
    elif has_docx:
        print(f"  • Word COM:        {'✓' if COMTYPES_AVAILABLE else '✗'} (Windows, best quality)")
        print(f"  • python-docx:     {'✓' if DOCX_AVAILABLE else '✗'} (for repair)")
        print(f"  • LibreOffice:     checking...")
    
    if not (COMTYPES_AVAILABLE or ASPOSE_AVAILABLE):
        print(f"\n{Fore.YELLOW}⚠ No primary conversion methods available!{Style.RESET_ALL}")
        print(f"{Fore.YELLOW}  Install: pip install comtypes (for best quality){Style.RESET_ALL}")
        if has_pptx:
            print(f"{Fore.YELLOW}  Or: pip install aspose-slides (cross-platform for PPTX){Style.RESET_ALL}")
        print(f"{Fore.YELLOW}  Falling back to LibreOffice if available...{Style.RESET_ALL}\n")
    
    print(f"\nFound {len(office_files)} Office files to convert")
    if has_pptx and has_docx:
        pptx_count = sum(1 for f in office_files if f.suffix.lower() in ['.pptx', '.ppt'])
        docx_count = sum(1 for f in office_files if f.suffix.lower() in ['.docx', '.doc'])
        print(f"  • PowerPoint: {pptx_count} file(s)")
        print(f"  • Word: {docx_count} file(s)")
    print()
    
    converter = OfficeConverter()
    converted_files = []
    failed_files = []
    stats = {'success': 0, 'repaired': 0, 'failed': 0}
    
    for idx, office_file in enumerate(office_files, 1):
        print(f"  [{idx}/{len(office_files)}] {office_file.name}")
        pdf_file = office_file.with_suffix(".pdf")
        
        # Determine file type and convert
        ext = office_file.suffix.lower()
        success = False
        method = "none"
        
        try:
            if ext in ['.pptx', '.ppt']:
                success, method = converter.convert_pptx_to_pdf(office_file, pdf_file)
            elif ext in ['.docx', '.doc']:
                success, method = converter.convert_docx_to_pdf(office_file, pdf_file)
            
            if success:
                size = pdf_file.stat().st_size
                print(f"    {Fore.GREEN}✓{Style.RESET_ALL} Converted using {method} ({size:,} bytes)")
                converted_files.append(pdf_file)
                
                if "repaired" in method:
                    stats['repaired'] += 1
                else:
                    stats['success'] += 1
                
                # Clean up source file after successful conversion
                try:
                    office_file.unlink()
                except:
                    pass
            else:
                print(f"    {Fore.RED}✗{Style.RESET_ALL} Failed - no conversion method succeeded")
                failed_files.append(office_file.name)
                stats['failed'] += 1
        
        except Exception as e:
            print(f"    {Fore.RED}✗{Style.RESET_ALL} Error: {str(e)[:60]}")
            failed_files.append(office_file.name)
            stats['failed'] += 1
    
    # Summary
    print(f"\n{Fore.GREEN}{'='*70}{Style.RESET_ALL}")
    print(f"{Fore.GREEN}✓ Conversion complete:{Style.RESET_ALL}")
    print(f"  Total files:           {len(office_files)}")
    print(f"  Successful:            {stats['success']}")
    print(f"  Repaired + converted:  {stats['repaired']}")
    print(f"  Failed:                {stats['failed']}")
    
    if stats['success'] + stats['repaired'] > 0:
        success_rate = ((stats['success'] + stats['repaired']) / len(office_files)) * 100
        print(f"  Success rate:          {success_rate:.1f}%")
    
    if failed_files:
        print(f"\n{Fore.YELLOW}Failed files:{Style.RESET_ALL}")
        for fname in failed_files[:10]:  # Show first 10
            print(f"  • {fname}")
        if len(failed_files) > 10:
            print(f"  ... and {len(failed_files) - 10} more")
    
    print(f"{Fore.GREEN}{'='*70}{Style.RESET_ALL}")
    
    return converted_files


def natural_sort_key(path: Path) -> list:
    """
    Generate a sort key for natural (numeric) sorting of filenames.
    Handles filenames like: 1.Name.pdf, 2.Name.pdf, 10.Name.pdf, 11.Name.pdf
    """
    import re
    parts = re.split(r'(\d+)', path.name)
    # Convert numeric parts to integers for proper comparison
    return [int(part) if part.isdigit() else part.lower() for part in parts]


def merge_pdfs_by_type(base_dir: Path, resource_types: List[str]):
    """Merge PDFs by resource type for each unit"""
    print(f"\n{Fore.CYAN}[6/7] Merging PDFs...{Style.RESET_ALL}")

    # Find all unit directories
    unit_dirs = sorted([d for d in base_dir.iterdir() if d.is_dir() and d.name.startswith("Unit_")])

    for unit_dir in unit_dirs:
        print(f"\n{Fore.BLUE}Processing {unit_dir.name}{Style.RESET_ALL}")

        # Merge by resource type
        for res_id in resource_types:
            resource_name = RESOURCE_TYPES[res_id]
            resource_dir = unit_dir / resource_name
            
            if resource_dir.exists():
                # Use natural sorting to preserve numeric order (1, 2, ..., 10, 11, not 1, 10, 11, 2)
                pdf_files = sorted(resource_dir.glob("*.pdf"), key=natural_sort_key)
                
                if pdf_files:
                    print(f"  {resource_name}: {len(pdf_files)} PDFs")

                    # Create merged PDF
                    output_file = unit_dir / f"{unit_dir.name}_{resource_name}_Merged.pdf"

                    try:
                        merger = PdfWriter()
                        for pdf_file in pdf_files:
                            try:
                                merger.append(str(pdf_file))
                                print(f"    + {pdf_file.name}")
                            except Exception as e:
                                print(f"    ✗ Failed to add {pdf_file.name}: {e}")

                        if len(merger.pages) > 0:
                            with open(output_file, "wb") as f:
                                merger.write(f)
                            merger.close()

                            size = output_file.stat().st_size
                            print(f"    {Fore.GREEN}✓ Created {output_file.name} ({size:,} bytes, {len(merger.pages)} pages){Style.RESET_ALL}")
                        else:
                            print(f"    {Fore.YELLOW}⚠ No valid PDFs to merge{Style.RESET_ALL}")

                    except Exception as e:
                        print(f"    {Fore.RED}✗ Merge failed: {e}{Style.RESET_ALL}")

    print(f"\n{Fore.GREEN}✓ PDF merging complete{Style.RESET_ALL}")


def cleanup_unwanted_files(base_dir: Path):
    """Remove README files and other unwanted files"""
    print(f"\n{Fore.CYAN}[7/7] Cleaning up unwanted files...{Style.RESET_ALL}")

    removed_count = 0
    
    # Extended list of unwanted file patterns
    unwanted_patterns = [
        "README*",
        "*.md",
        "*.txt",
        "Thumbs.db",
        ".DS_Store",
        "desktop.ini",
        "*.tmp",
        "*.temp",
        "*~",  # Backup files
    ]
    
    # Also remove empty directories
    for pattern in unwanted_patterns:
        for file in base_dir.rglob(pattern):
            if file.is_file():
                try:
                    file.unlink()
                    print(f"  {Fore.YELLOW}✓{Style.RESET_ALL} Removed: {file.name}")
                    removed_count += 1
                except Exception as e:
                    print(f"  {Fore.RED}✗{Style.RESET_ALL} Failed to remove {file.name}: {e}")

    # Remove empty directories
    empty_dirs = []
    for dirpath in base_dir.rglob("*"):
        if dirpath.is_dir():
            try:
                # Check if directory is empty
                if not any(dirpath.iterdir()):
                    empty_dirs.append(dirpath)
            except:
                pass
    
    # Remove empty directories (deepest first)
    for empty_dir in sorted(empty_dirs, key=lambda p: len(p.parts), reverse=True):
        try:
            empty_dir.rmdir()
            print(f"  {Fore.YELLOW}✓{Style.RESET_ALL} Removed empty directory: {empty_dir.name}")
            removed_count += 1
        except Exception as e:
            pass  # Directory might not be empty anymore

    if removed_count > 0:
        print(f"{Fore.GREEN}✓ Removed {removed_count} unwanted items{Style.RESET_ALL}")
    else:
        print(f"{Fore.GREEN}✓ No unwanted files found{Style.RESET_ALL}")


def display_courses(courses: List[Dict], page_size: int = 10, fetch_fn=None) -> Optional[Dict]:
    """
    Interactive live-search course selector using curses.
    Shows a centered search box; list filters dynamically as the user types.
    Returns the selected course dict, or None if user pressed Esc/q.
    
    Args:
        courses: List of course dictionaries
        page_size: Number of results to show at once (default 10)
        fetch_fn: Optional callable to fetch fresh courses from server (called on F5)
    
    Note: On Windows, requires 'windows-curses' package:
          pip install windows-curses
    """
    if not courses:
        print("No courses available.")
        return None

    result = [None]  # use list so inner function can write to it

    def _ui(stdscr):
        curses.curs_set(1)          # show cursor
        curses.use_default_colors()
        curses.init_pair(1, curses.COLOR_BLACK, curses.COLOR_CYAN)   # highlight
        curses.init_pair(2, curses.COLOR_CYAN,  -1)                  # header / border
        curses.init_pair(3, curses.COLOR_YELLOW, -1)                 # course code
        curses.init_pair(4, curses.COLOR_GREEN,  -1)                 # footer hint

        query        = ""
        selected_idx = 0    # index in the *filtered* list
        scroll_top   = 0    # first visible row in the filtered list
        LIST_SIZE    = 10   # max rows shown at once

        def matches(course: Dict, q: str) -> bool:
            q = q.lower()
            return (q in course.get("subjectCode", "").lower() or
                    q in course.get("subjectName", "").lower())

        while True:
            stdscr.erase()
            max_h, max_w = stdscr.getmaxyx()

            # ── filtered list ─────────────────────────────────────────────
            filtered = [c for c in courses if matches(c, query)] if query else courses[:]

            # clamp selection
            if selected_idx >= len(filtered):
                selected_idx = max(0, len(filtered) - 1)
            if selected_idx < scroll_top:
                scroll_top = selected_idx
            if selected_idx >= scroll_top + LIST_SIZE:
                scroll_top = selected_idx - LIST_SIZE + 1

            # ── box dimensions ────────────────────────────────────────────
            box_w   = min(max_w - 4, 78)
            # header(1) + border-top(1) + search(1) + divider(1) +
            # list rows + count(1) + hint(1) + border-bot(1) = LIST_SIZE + 7
            box_h   = LIST_SIZE + 7
            box_x   = (max_w - box_w) // 2
            box_y   = max(0, (max_h - box_h) // 2)

            inner_w = box_w - 2   # width inside the border

            # ── draw outer border ─────────────────────────────────────────
            try:
                # top
                stdscr.addstr(box_y,         box_x, "╔" + "═" * (box_w - 2) + "╗", curses.color_pair(2))
                # bottom
                stdscr.addstr(box_y + box_h, box_x, "╚" + "═" * (box_w - 2) + "╝", curses.color_pair(2))
                # sides
                for row in range(1, box_h):
                    stdscr.addstr(box_y + row, box_x,           "║", curses.color_pair(2))
                    stdscr.addstr(box_y + row, box_x + box_w - 1, "║", curses.color_pair(2))
            except curses.error:
                pass

            # ── header ────────────────────────────────────────────────────
            header = " PESU Academy — Select Course "
            header = header[:inner_w]
            hx = box_x + 1 + (inner_w - len(header)) // 2
            try:
                stdscr.addstr(box_y + 1, hx, header, curses.color_pair(2) | curses.A_BOLD)
            except curses.error:
                pass

            # ── divider after header ──────────────────────────────────────
            try:
                stdscr.addstr(box_y + 2, box_x, "╠" + "═" * (box_w - 2) + "╣", curses.color_pair(2))
            except curses.error:
                pass

            # ── search bar ────────────────────────────────────────────────
            prompt = " 🔍 "
            search_label = (prompt + query)[:inner_w]
            try:
                stdscr.addstr(box_y + 3, box_x + 1, search_label.ljust(inner_w))
            except curses.error:
                pass

            # ── divider after search ──────────────────────────────────────
            try:
                stdscr.addstr(box_y + 4, box_x, "╠" + "─" * (box_w - 2) + "╣", curses.color_pair(2))
            except curses.error:
                pass

            # ── course list ───────────────────────────────────────────────
            visible = filtered[scroll_top: scroll_top + LIST_SIZE]
            for i, course in enumerate(visible):
                abs_i     = scroll_top + i
                row_y     = box_y + 5 + i
                is_sel    = abs_i == selected_idx

                code = course.get("subjectCode", "")[:18]
                name = course.get("subjectName", "")
                # trim name so line fits inside box
                max_name = inner_w - len(code) - 3
                if len(name) > max_name:
                    name = name[:max_name - 1] + "…"

                line = f" {code:<18}  {name}"
                line = line[:inner_w]

                try:
                    if is_sel:
                        stdscr.addstr(row_y, box_x + 1, line.ljust(inner_w), curses.color_pair(1) | curses.A_BOLD)
                    else:
                        stdscr.addstr(row_y, box_x + 1, f" {code:<18}", curses.color_pair(3))
                        rest = f"  {name}"
                        stdscr.addstr(row_y, box_x + 1 + 19, rest[:inner_w - 19])
                except curses.error:
                    pass

            # fill blank rows if fewer than LIST_SIZE results
            for i in range(len(visible), LIST_SIZE):
                try:
                    stdscr.addstr(box_y + 5 + i, box_x + 1, " " * inner_w)
                except curses.error:
                    pass

            # ── show "no results" message if search yields nothing ───────
            if len(filtered) == 0 and query:
                msg = " No results. Press Ctrl+R to fetch from server, Esc to cancel "
                try:
                    stdscr.addstr(box_y + 5 + LIST_SIZE // 2, box_x + 1,
                                  msg[:inner_w].center(inner_w),
                                  curses.color_pair(3) | curses.A_BOLD)
                except curses.error:
                    pass

            # ── count row ─────────────────────────────────────────────────
            count_row = box_y + 5 + LIST_SIZE
            count_str = f" {len(filtered)} course(s) found"
            try:
                stdscr.addstr(count_row, box_x + 1, count_str[:inner_w].ljust(inner_w), curses.color_pair(2))
            except curses.error:
                pass

            # ── footer hint ───────────────────────────────────────────────
            hint = " ↑↓ navigate   Enter select   Esc/q quit "
            hint = hint[:inner_w]
            try:
                stdscr.addstr(box_y + box_h - 1, box_x + 1, hint.ljust(inner_w), curses.color_pair(4))
            except curses.error:
                pass

            # ── position cursor inside search bar ────────────────────────
            cursor_x = box_x + 1 + len(prompt) + len(query)
            cursor_x = min(cursor_x, box_x + box_w - 2)
            try:
                stdscr.move(box_y + 3, cursor_x)
            except curses.error:
                pass

            stdscr.refresh()

            # ── handle input ──────────────────────────────────────────────
            try:
                key = stdscr.get_wch()
            except curses.error:
                continue

            if isinstance(key, str):
                # printable character → append to query
                if key in ("\x1b",):          # Esc
                    result[0] = None
                    return
                elif key in ("\n", "\r"):      # Enter
                    if filtered:
                        result[0] = filtered[selected_idx]
                    return
                elif key in ("\x7f", "\b"):   # Backspace
                    query = query[:-1]
                    selected_idx = 0
                    scroll_top   = 0
                elif key == "\x12" and fetch_fn is not None:  # Ctrl+R - Refresh
                    # Exit curses temporarily, fetch, re-enter
                    curses.endwin()
                    print(f"\n{Fore.CYAN}Fetching fresh course list from server...{Style.RESET_ALL}")
                    try:
                        fresh_courses = fetch_fn()
                        courses.clear()
                        courses.extend(fresh_courses)
                        selected_idx = 0
                        scroll_top = 0
                        query = ""  # Reset search after refresh
                        print(f"{Fore.GREEN}✓ Done! {len(courses)} courses loaded. Press any key to continue...{Style.RESET_ALL}")
                        input()
                    except Exception as e:
                        print(f"{Fore.RED}✗ Fetch failed: {e}{Style.RESET_ALL}")
                        input("Press Enter to go back...")
                    stdscr.refresh()
                elif key == "q" and not query:  # q only quits when search is empty
                    result[0] = None
                    return
                elif key.isprintable():
                    query += key
                    selected_idx = 0
                    scroll_top   = 0

            else:  # special key (integer)
                if key == curses.KEY_UP:
                    if selected_idx > 0:
                        selected_idx -= 1
                        if selected_idx < scroll_top:
                            scroll_top = selected_idx
                elif key == curses.KEY_DOWN:
                    if selected_idx < len(filtered) - 1:
                        selected_idx += 1
                        if selected_idx >= scroll_top + LIST_SIZE:
                            scroll_top += 1
                elif key == curses.KEY_BACKSPACE:
                    query = query[:-1]
                    selected_idx = 0
                    scroll_top   = 0
                elif key in (curses.KEY_ENTER, 10, 13):
                    if filtered:
                        result[0] = filtered[selected_idx]
                    return

    try:
        curses.wrapper(_ui)
    except Exception as e:
        # Fallback in case curses fails (e.g., not in a terminal)
        print(f"\n{Fore.YELLOW}⚠ Curses UI not available: {e}{Style.RESET_ALL}")
        print(f"{Fore.YELLOW}  On Windows, install: pip install windows-curses{Style.RESET_ALL}")
        print(f"{Fore.YELLOW}  Falling back to basic selection...{Style.RESET_ALL}\n")
        
        # Simple fallback: just list and pick by number
        for i, course in enumerate(courses, 1):
            code = course.get('subjectCode', 'N/A')
            name = course.get('subjectName', 'N/A')
            print(f"{i:3d}. {code:<20} {name}")
        
        try:
            choice = input(f"\n{Fore.CYAN}Enter course number (or 'q' to quit): {Style.RESET_ALL}").strip()
            if choice.lower() == 'q':
                return None
            idx = int(choice) - 1
            if 0 <= idx < len(courses):
                return courses[idx]
        except (ValueError, IndexError):
            print(f"{Fore.RED}Invalid selection{Style.RESET_ALL}")
            return None
    
    return result[0]


def main():
    """Main interactive flow"""
    print(f"\n{Fore.GREEN}{Style.BRIGHT}{'='*80}{Style.RESET_ALL}")
    print(f"{Fore.GREEN}{Style.BRIGHT}  PESU Academy Interactive Downloader{Style.RESET_ALL}")
    print(f"{Fore.GREEN}{Style.BRIGHT}{'='*80}{Style.RESET_ALL}\n")

    # Get credentials
    try:
        from dotenv import load_dotenv

        load_dotenv()
        username = os.getenv("PESU_USERNAME")
        password = os.getenv("PESU_PASSWORD")
    except:
        username = None
        password = None

    if not username or not password:
        print(f"{Fore.CYAN}Enter your PESU Academy credentials:{Style.RESET_ALL}")
        username = input("Username (SRN): ").strip()
        password = getpass.getpass("Password: ").strip()

    if not username or not password:
        print(f"{Fore.RED}Username and password are required!{Style.RESET_ALL}")
        sys.exit(1)

    # Initialize downloader and login
    downloader = PESUInteractiveDownloader(username, password)

    try:
        downloader.login()

        # Get courses
        courses = downloader.get_courses()

        # Filter by academic year
        courses = downloader.filter_courses_by_year(courses)

        if not courses:
            print(f"\n{Fore.RED}No courses found for selected year. Exiting.{Style.RESET_ALL}")
            return

        # Display and select course
        print(f"\n{Fore.CYAN}[3/7] Selecting course...{Style.RESET_ALL}")
        selected_course = display_courses(courses, fetch_fn=downloader._fetch_and_cache_courses)

        if not selected_course:
            print(f"\n{Fore.YELLOW}No course selected. Exiting.{Style.RESET_ALL}")
            return

        print(
            f"\n{Fore.GREEN}✓ Selected: {selected_course['subjectCode']} - {selected_course['subjectName']}{Style.RESET_ALL}"
        )

        course_id = selected_course["id"]
        course_code = selected_course["subjectCode"]
        course_name = selected_course["subjectName"]

        # Get units
        all_units = downloader.get_units(course_id)
        print(f"\n{Fore.CYAN}Available Units:{Style.RESET_ALL}")
        for idx, unit in enumerate(all_units, 1):
            print(f"  {idx}. {unit['name']}")

        # Select units
        print(
            f"\n{Fore.CYAN}Enter unit numbers (e.g., 1,2,3 for units 1,2,3 or 'all'): {Style.RESET_ALL}",
            end="",
        )
        unit_input = input().strip().lower()

        if unit_input == "all":
            selected_units = list(range(1, len(all_units) + 1))
        else:
            selected_units = [int(x.strip()) for x in unit_input.split(",")]

        print(f"{Fore.GREEN}✓ Selected units: {', '.join(map(str, selected_units))}{Style.RESET_ALL}")

        # Select resource types
        print(f"\n{Fore.CYAN}Available Resource Types:{Style.RESET_ALL}")
        for res_id, res_name in RESOURCE_TYPES.items():
            print(f"  {res_id}. {res_name}")

        print(
            f"\n{Fore.CYAN}Enter resource type numbers (e.g., 2,3,6 for Slides, Notes, QB or 'all'): {Style.RESET_ALL}",
            end="",
        )
        resource_input = input().strip().lower()

        if resource_input == "all":
            selected_resources = list(RESOURCE_TYPES.keys())
        else:
            selected_resources = [x.strip() for x in resource_input.split(",")]

        resource_names = [RESOURCE_TYPES[r] for r in selected_resources]
        print(f"{Fore.GREEN}✓ Selected resources: {', '.join(resource_names)}{Style.RESET_ALL}")

        # Create base directory
        base_dir = Path("downloads") / f"{course_code}"
        base_dir.mkdir(parents=True, exist_ok=True)

        # Download resources
        downloader.download_resources(
            course_id, course_name, selected_units, selected_resources, base_dir
        )

        # Check for non-PDF files and ask for conversion
        office_files = []
        for ext in ["*.docx", "*.pptx", "*.doc", "*.ppt"]:
            office_files.extend(base_dir.rglob(ext))

        if office_files:
            print(
                f"\n{Fore.YELLOW}Found {len(office_files)} non-PDF files (Word/PowerPoint){Style.RESET_ALL}"
            )
            print(
                f"{Fore.CYAN}Do you want to convert them to PDF? (y/n): {Style.RESET_ALL}",
                end="",
            )
            convert_choice = input().strip().lower()

            if convert_choice == "y":
                convert_office_to_pdf(base_dir)

        # Ask for PDF merging
        print(f"\n{Fore.CYAN}Do you want to merge PDFs by resource type? (y/n): {Style.RESET_ALL}", end="")
        merge_choice = input().strip().lower()

        if merge_choice == "y":
            merge_pdfs_by_type(base_dir, selected_resources)

        # Automatically cleanup unwanted files
        cleanup_unwanted_files(base_dir)

        # Final summary
        print(f"\n{Fore.GREEN}{Style.BRIGHT}{'='*80}{Style.RESET_ALL}")
        print(f"{Fore.GREEN}{Style.BRIGHT}✓ All tasks completed!{Style.RESET_ALL}")
        print(f"{Fore.GREEN}Location: {base_dir.absolute()}{Style.RESET_ALL}")
        print(f"{Fore.GREEN}{Style.BRIGHT}{'='*80}{Style.RESET_ALL}\n")

    except Exception as e:
        print(f"\n{Fore.RED}Error: {e}{Style.RESET_ALL}")
        sys.exit(1)


if __name__ == "__main__":
    main()
    